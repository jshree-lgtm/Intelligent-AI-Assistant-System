import streamlit as st
import os
import pickle
import faiss
import numpy as np

from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# Embedding model
embedder = SentenceTransformer("all-MiniLM-L6-v2")

VECTOR_DB_PATH = "vector_db/faiss_index.bin"
METADATA_PATH = "vector_db/metadata.pkl"


def chunk_text(text, chunk_size=500):

    chunks = []

    for i in range(0, len(text), chunk_size):
        chunks.append(text[i:i + chunk_size])

    return chunks


def read_txt(file):
    return file.read().decode("utf-8")


def read_pdf(file):

    pdf = PdfReader(file)

    text = ""

    for page in pdf.pages:

        page_text = page.extract_text()

        if page_text:
            text += page_text + "\n"

    return text


def read_docx(file):

    doc = Document(file)

    text = ""

    for para in doc.paragraphs:
        text += para.text + "\n"

    return text


def read_pptx(file):

    prs = Presentation(file)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def show():

    st.header("📚 Automatic Knowledge Base Update")

    if st.button("🗑 Reset Knowledge Base"):

        if os.path.exists(VECTOR_DB_PATH):
            os.remove(VECTOR_DB_PATH)

        if os.path.exists(METADATA_PATH):
            os.remove(METADATA_PATH)

        st.success("Knowledge base has been reset successfully!")
        st.rerun()

    st.write(
        "Upload TXT, PDF, DOCX, or PPTX files to expand the chatbot's knowledge base."
    )

    uploaded_file = st.file_uploader(
        "Upload Knowledge Document",
        type=["txt", "pdf", "docx", "pptx"]
    )

    if uploaded_file is not None:

        try:

            file_type = uploaded_file.name.split(".")[-1].lower()

            if file_type == "txt":
                text = read_txt(uploaded_file)

            elif file_type == "pdf":
                text = read_pdf(uploaded_file)

            elif file_type == "docx":
                text = read_docx(uploaded_file)

            elif file_type == "pptx":
                text = read_pptx(uploaded_file)

            else:
                st.error("Unsupported file type")
                return

            if len(text.strip()) == 0:
                st.warning("No readable text found in the document.")
                return

            chunks = chunk_text(text)

            embeddings = embedder.encode(chunks)

            embeddings = np.array(
                embeddings,
                dtype="float32"
            )

            dimension = embeddings.shape[1]

            if os.path.exists(VECTOR_DB_PATH):

                index = faiss.read_index(
                    VECTOR_DB_PATH
                )

                with open(
                    METADATA_PATH,
                    "rb"
                ) as f:

                    metadata = pickle.load(f)

            else:

                index = faiss.IndexFlatL2(
                    dimension
                )

                metadata = []

            index.add(embeddings)

            metadata.extend(chunks)

            faiss.write_index(
                index,
                VECTOR_DB_PATH
            )

            with open(
                METADATA_PATH,
                "wb"
            ) as f:

                pickle.dump(
                    metadata,
                    f
                )

            st.success(
                f"Knowledge Base Updated! Added {len(chunks)} chunks."
            )

        except Exception as e:

            st.error(
                f"Error processing document: {e}"
            )

    st.subheader("📊 Vector Database Status")

    if os.path.exists(METADATA_PATH):

        with open(
            METADATA_PATH,
            "rb"
        ) as f:

            metadata = pickle.load(f)

        st.write(
            f"Total Knowledge Chunks: {len(metadata)}"
        )

    else:

        st.write(
            "Vector database is empty."
        )